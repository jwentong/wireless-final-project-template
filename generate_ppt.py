"""Generate the defense slide deck (.pptx).

Narrative follows the course's Vibe -> Spec -> Harness engineering workflow, with
the mock-driven design revision and the implementation debugging as highlights.
Run: python generate_ppt.py  ->  答辩PPT.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

RESULTS = Path("results")
OUTPUT = Path("答辩PPT.pptx")

TEAL = RGBColor(0x2A, 0x6F, 0x78)
ORANGE = RGBColor(0xC2, 0x69, 0x3F)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set(frame, size, color, bold=False):
    for para in frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Microsoft YaHei"


def title_bar(slide, text, idx=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    tb = bar.text_frame
    tb.margin_left = Inches(0.5)
    tb.word_wrap = True
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb.text = text
    _set(tb, 26, WHITE, bold=True)
    if idx is not None:
        n = slide.shapes.add_textbox(Inches(12.3), Inches(0.35), Inches(0.9), Inches(0.5))
        n.text_frame.text = idx
        _set(n.text_frame, 14, WHITE)


def bullets(slide, items, left=0.6, top=1.5, width=12.1, height=5.6, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ("• " if level == 0 else "   – ") + text
        para.level = level
        para.space_after = Pt(8)
        for run in para.runs:
            run.font.size = Pt(size - level * 2)
            run.font.color.rgb = DARK if level == 0 else GRAY
            run.font.name = "Microsoft YaHei"
            if level == 0 and text.endswith("："):
                run.font.bold = True
    return box


def picture(slide, name, left, top, width):
    path = RESULTS / name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def new(idx=None, title=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        title_bar(s, title, idx)
    return s


# ---------------- 1 Cover ----------------
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(1, Inches(0), Inches(2.4), Inches(13.333), Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = TEAL; band.line.fill.background()
t = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.5))
t.text_frame.word_wrap = True
t.text_frame.text = "基于 AI 辅助编程的无线通信文件传输基带仿真系统"
_set(t.text_frame, 34, WHITE, bold=True)
sub = s.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.7))
sub.text_frame.text = "无线通信技术期末项目答辩 · Vibe → Spec → Harness 工程流程"
_set(sub.text_frame, 18, WHITE)
info = s.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2))
info.text_frame.text = "姓名：齐宇恒    学号：2024140044    GitHub：EchoTreee"
info.text_frame.add_paragraph().text = "Fork：github.com/EchoTreee/wireless-final-project-template"
_set(info.text_frame, 16, GRAY)

# ---------------- 2 Overview ----------------
s = new("02", "项目概述：一句话与一条链路")
bullets(s, [
    ("目标：把 Test.txt（UTF-8 文本）经完整无线基带链路传输，在 received.txt 无损恢复", 0),
    ("固定链路：源编码 → 扰码 → 信道编码 → 帧 → QPSK → 信道 → 同步 → 解调 → 译码 → 解扰 → 源解码", 0),
    ("验收基线：SNR ≥ 12 dB、AWGN、固定 seed 下 received.txt 与 Test.txt 完全一致", 0),
    ("当前成果：", 0),
    ("单元测试 47 项、公开测试 22 项全部通过", 1),
    ("端到端 BER=0、text_match_rate=1.0、checksum_pass=true", 1),
    ("完成 Level 3：卷积码 Viterbi、衰落+均衡、多调制对比", 1),
])

# ---------------- 3 Strategy ----------------
s = new("03", "我的策略：满分点不只在代码")
bullets(s, [
    ("评分结构：设计20 + mock15 + 实现25 + 测试20 + 分析10 + AI与答辩10", 0),
    ("洞察：其中约 45 分落在‘工程过程 + 可解释性’，而非纯代码", 0),
    ("策略：", 0),
    ("真实走 PRD→DESIGN→TEST_PLAN→mock→TDD→验证 流程并全程留痕", 1),
    ("做到 Level 3 提高模块", 1),
    ("每个决策都记进 AI_LOG，既拿分又是答辩故事素材", 1),
])

# ---------------- 4 Workflow ----------------
s = new("04", "工程流程：Vibe → Spec → Harness")
bullets(s, [
    ("Vibe（探索）：最低成本澄清需求，读 PRD 与公开测试，定满分策略", 0),
    ("Spec（收敛）：把隐含判断变成显式约束 → DESIGN.md / TEST_PLAN.md", 0),
    ("Harness（验证）：mock 验证设计 → TDD 实现 → 单元/公开测试 → 完成前验证", 0),
    ("关键习惯：先看 public_tests 反推接口契约，让测试成为设计的 source of truth", 0),
], size=19)

# ---------------- 5 Architecture ----------------
s = new("05", "系统架构：模块化与固定链路")
bullets(s, [
    ("src/ 模块与公开测试自动发现命名一致，零适配通过：", 0),
    ("source / scramble / channel_coding / framing / modulation", 1),
    ("channel / synchronization / equalizer / metrics / pipeline", 1),
    ("main.py：统一 CLI，非交互，输出 received.txt + metrics.json + 三张图", 0),
    ("可复现：所有随机过程由 seed 控制；通用性：对任意 UTF-8 / 长度 / SNR 工作，无硬编码", 0),
])

# ---------------- 6 Design: frame ----------------
s = new("06", "关键设计①：帧结构")
bullets(s, [
    ("帧 = Barker-13 前导 | orig_len×3 | coded_len×3 | payload | CRC-16 | padding", 0),
    ("length 双字段：orig_len（原始 bit 数，去 padding 恢复文本）+ coded_len（编码后 bit 数，定位 payload）", 0),
    ("解决了 PRD‘length 语义’与‘接收端定位’的冲突", 1),
    ("前导 Barker-13：QPSK 后落主对角线，自相关旁瓣≤1，同步峰尖锐", 0),
], top=1.4, height=3.0)
picture(s, "sync_peak.png", 3.2, 4.3, 7.0)

# ---------------- 7 Design: QPSK + SNR ----------------
s = new("07", "关键设计②：QPSK 与 SNR 口径")
bullets(s, [
    ("Gray 映射 QPSK：00→(1+j)/√2，01→(-1+j)/√2，11→(-1-j)/√2，10→(1-j)/√2", 0),
    ("相邻象限仅差 1 bit；1/√2 归一化使符号平均功率=1", 1),
    ("SNR 定义：符号平均功率 / 复噪声平均功率（dB），与 PRD 一致", 0),
    ("Eb/N0 = SNR − 3.01 dB（QPSK 每符号 2 bit），用于对照理论 BER", 1),
], top=1.4, height=3.0)
picture(s, "constellation.png", 4.4, 4.2, 4.6)

# ---------------- 8 Design: conv code ----------------
s = new("08", "关键设计③：卷积码 + Viterbi")
bullets(s, [
    ("卷积码 K=7、码率 1/2、生成多项式 (171,133)₈，发送端加 6 零尾比特终止", 0),
    ("接收端硬判决 Viterbi：64 状态网格加比选 + 回溯最优路径", 0),
    ("零尾终止保证无噪声完全可逆；相比无编码约 4–5 dB 编码增益", 0),
    ("同时实现 Hamming(7,4) 作对比基线", 0),
])

# ---------------- 9 Mock story ----------------
s = new("09", "亮点①：mock 驱动的设计修订")
bullets(s, [
    ("先 mock 后实现：用重复码占位搭骨架，压低 SNR 扫描鲁棒性", 0),
    ("发现缺陷：PRD 链路中信道编码在帧封装之前 → 帧头 length 字段不受 FEC 保护", 0),
    ("低 SNR 单比特错落在 32-bit 长度字段 → 整帧解析崩溃", 1),
    ("修订：length 字段 3× 重复 + 多数判决（不改链路顺序，+128 bit）", 0),
], top=1.4, height=3.2)
rows = [("SNR", "修订前 header 损坏", "修订后 header 损坏"),
        ("4 dB", "18/20", "6/20"), ("6 dB", "11/20", "1/20")]
tbl = s.shapes.add_table(len(rows), 3, Inches(3.3), Inches(4.7), Inches(6.7), Inches(2.0)).table
for i, r in enumerate(rows):
    for j, v in enumerate(r):
        tbl.cell(i, j).text = v
        for para in tbl.cell(i, j).text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(14); run.font.name = "Microsoft YaHei"
                if i == 0:
                    run.font.bold = True

# ---------------- 10 Impl debugging ----------------
s = new("10", "亮点②：实现期两个调试")
bullets(s, [
    ("调试一 · checksum 假阴性：CRC 原本算在 coded payload 上", 0),
    ("AWGN 12dB 残余比特错被 Viterbi 纠正→文本对，但译码前 CRC 失败，矛盾", 1),
    ("修复：CRC 覆盖原始信息位、随信息进 FEC 受保护，译码后校验通过（符合 PRD）", 1),
    ("调试二 · MMSE 的 SNR 定义：单测噪声功率与全系统 SNR 约定不一致", 0),
    ("导致 MMSE 过度收缩、MSE 反大于 ZF；统一约定后 MMSE 恢复最优", 1),
    ("两处都定位真实根因（协议分层 / SNR 定义），不是改测试妥协", 0),
], size=17)

# ---------------- 11 Test results ----------------
s = new("11", "测试结果：全绿 + 端到端无损")
bullets(s, [
    ("单元测试 47 项全绿（含卷积码纠错、framing header 纠错、AWGN 可复现等）", 0),
    ("教师公开测试 22 项全绿（结构/文档/各模块/端到端/反硬编码）", 0),
    ("端到端 @SNR12 seed2026：received.txt == Test.txt，BER=0，match=1.0，checksum=true", 0),
    ("鲁棒性：对不同文本/长度/SNR/seed/0–128 符号同步偏移/异常输入均处理", 0),
])

# ---------------- 12 Experiment: coding gain ----------------
s = new("12", "实验分析：BER-SNR 与编码增益")
picture(s, "fec_comparison.png", 0.7, 1.5, 6.0)
bullets(s, [
    ("卷积码 < 汉明码 < 无编码", 0),
    ("编码增益约 4–5 dB", 0),
    ("10 dB：无编码 BER≈5e-4，卷积码已无误码", 0),
    ("星座图随 SNR 扩散越界即产生误比特", 0),
], left=7.0, top=1.8, width=5.8, height=4.5, size=17)

# ---------------- 13 Experiment: fading ----------------
s = new("13", "提高模块：衰落信道与均衡")
picture(s, "fading_comparison.png", 0.7, 1.5, 6.0)
bullets(s, [
    ("Rayleigh 无均衡：BER 居高（随机旋转/缩放）", 0),
    ("均衡后 BER 随 SNR 下降", 0),
    ("ZF ≈ MMSE for QPSK：仅相位判决，差一个正实标量 → 硬判决 BER 重合", 0),
    ("MMSE 在估计 MSE 意义下更优（单测已验证）", 0),
], left=7.0, top=1.8, width=5.8, height=4.5, size=16)

# ---------------- 14 Experiment: modulation ----------------
s = new("14", "提高模块：多调制对比与自适应")
picture(s, "modulation_comparison.png", 0.7, 1.5, 6.0)
bullets(s, [
    ("BPSK < QPSK < 16-QAM（符号 SNR 口径）", 0),
    ("QPSK 较 BPSK 差约 3 dB：每符号 2 bit，Eb/N0 低 3 dB", 0),
    ("自适应依据：高 SNR 选高阶提速，低 SNR 退低阶保可靠", 0),
], left=7.0, top=1.8, width=5.8, height=4.5, size=17)

# ---------------- 15 AI + integrity ----------------
s = new("15", "AI 协作与学术诚信")
bullets(s, [
    ("全程 AI_LOG：关键 prompt、AI 生成、人工修改、采纳理由", 0),
    ("保留：整体架构与接口（对齐公开测试）", 0),
    ("修改：length→32bit、帧头保护、CRC 分层、MMSE SNR 定义", 0),
    ("拒绝：‘直接生成最终代码’，坚持流程留痕", 0),
    ("全部代码可解释；无硬编码公开测试；无绕过链路的文件直拷", 0),
])

# ---------------- 16 Conclusion ----------------
s = new("16", "总结与答辩要点")
bullets(s, [
    ("达成：PRD 全口径达标 + Level 3 + 全测试绿 + 端到端无损", 0),
    ("最大亮点：mock 用数据发现并修复设计缺陷（header 保护）", 0),
    ("可深答的问题：", 0),
    ("为何 QPSK / 帧如何支持同步与纠错 / SNR 降低谁先失效 / 乱码排查顺序", 1),
    ("AI 生成内容保留、修改、拒绝了什么、为什么", 1),
])

prs.save(OUTPUT)
print(f"saved: {OUTPUT.resolve()}  slides={len(prs.slides._sldIdLst)}")
